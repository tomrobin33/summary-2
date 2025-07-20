import requests
import tempfile
import os
from pptx import Presentation
import pandas as pd
from docx import Document

def download_file(url):
    response = requests.get(url)
    response.raise_for_status()
    suffix = url.split('.')[-1].split('?')[0].lower()
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.' + suffix)
    tmp.write(response.content)
    tmp.close()
    return tmp.name, suffix

def parse_pptx(file_path):
    prs = Presentation(file_path)
    slides = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides.append({"slide": slide_text})
    return {"type": "ppt", "content": slides}

def parse_excel(file_path):
    xls = pd.ExcelFile(file_path)
    sheets = {}
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name)
        sheets[sheet_name] = df.fillna("").values.tolist()
    return {"type": "excel", "content": sheets}

def parse_word(file_path):
    doc = Document(file_path)
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    return {"type": "word", "content": paragraphs}

def parse_file(url):
    file_path, ext = download_file(url)
    try:
        if ext == "pptx":
            result = parse_pptx(file_path)
        elif ext in ["xlsx", "xls"]:
            result = parse_excel(file_path)
        elif ext == "docx":
            result = parse_word(file_path)
        else:
            result = {"error": "不支持的文件类型"}
    finally:
        os.remove(file_path)
    return result 